# -*- coding: utf-8 -*-
"""
Created on Sun Mar  1 09:38:50 2020
@author: Moojin Bin, Associate, Global Asset Allocation Strategy Department

"""


# %% Import Packages

from bs4 import BeautifulSoup
import pandas as pd
import numpy as np
import re
from pandas_datareader import data as pdr
import yfinance as yf
yf.pdr_override()
from datetime import datetime, timedelta
import datetime as dt
import matplotlib.pyplot as plt 
import matplotlib.ticker as tick
import matplotlib.gridspec as gridspec
import matplotlib.dates as mdates
from matplotlib.dates import DateFormatter
plt.rcParams.update({'font.size':15})
plt.rcParams.update({'figure.max_open_warning': 0})
plt.rcParams['axes.unicode_minus'] = False
import math
import seaborn as sns
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import random
from selenium import webdriver



# %% Date Setting
def make_beg_end_date(now):
    today = dt.date(now.year, now.month, now.day)    
    start_date = '1990-01-01'

    if today.strftime('%A') == 'Saturday':
        end_date = today - timedelta(days=1)    
    elif today.strftime('%A') == 'Sunday':
        end_date = today - timedelta(days=2)    
    else:
        end_date = today

    return start_date, end_date

now = datetime.now()
start_date, end_date = make_beg_end_date(now)

# %% Date Setting for tradibility
def make_beg_end_date2(now):
    today = dt.date(now.year, now.month, now.day)    
    start_date2 = today - timedelta(days=430)

    if today.strftime('%A') == 'Saturday':
        end_date2 = today - timedelta(days=1)    
    elif today.strftime('%A') == 'Sunday':
        end_date2 = today - timedelta(days=2)    
    else:
        end_date2 = today

    return start_date2, end_date2

now = datetime.now()
start_date2, end_date2 = make_beg_end_date2(now)


# %% Data From Excel
save_path = r'C:/ETF'
df_from_excel = pd.read_excel(save_path + '\ETF Description.xlsx',
                              index_col = 'symbol')
ETF_description = df_from_excel.copy().dropna()



# %% Crawling Functions

ticker = 'ticker'

def tag_to_list(tag):
    temp = tag.get_text().split('\n')
    return_list = list(filter(None, temp))
    return return_list

def list_to_df(input_list):
    df_columns = re.sub(r'\[[^)]*\]', '', input_list[0])
    df_index = input_list[1::2]
    df = pd.DataFrame(input_list[2::2], index = df_index, columns = [df_columns])
    df.index = df.index.str.replace(';','')
    return df

def make_top_holdings(l):
    holding = l[1::3]
    weight  = l[2::3]
    top_holdings = pd.DataFrame(weight, index=holding)
    top_holdings.columns = ['weight']
    return top_holdings

def make_soup1(ticker):
    source1 = "https://www.etf.com/"
    url1 = source1 + ticker
    driver = webdriver.Chrome('C:\ETF\chromedriver.exe')
    #driver.implicitly_wait(random.uniform(3, 5))
    driver.get(url1)
    html = driver.page_source
    driver.quit()
    soup1 = BeautifulSoup(html, 'html.parser')
    return soup1

def make_soup2(ticker):
    source2 = "https://www.etfdb.com/etf/"
    url2 = source2 + ticker + '/#holdings'
    driver = webdriver.Chrome('C:\ETF\chromedriver.exe')
    #driver.implicitly_wait(random.uniform(3, 5))
    driver.get(url2)
    html = driver.page_source
    driver.quit()
    soup2 = BeautifulSoup(html, 'html.parser')
    return soup2

def make_soup3(ticker):
    source3 = "https://www.etfdb.com/etf/"
    url3 = source3 + ticker + '/#charts'
    driver = webdriver.Chrome('C:\ETF\chromedriver.exe')
    driver.implicitly_wait(random.uniform(3, 5))
    driver.get(url3)
    html = driver.page_source
    driver.quit()
    soup2 = BeautifulSoup(html, 'html.parser')
    return soup2

def etf_name(soup):
    '''
    * soup = make_soup1(ticker)
    * This function will return the full name of ETF from ETF.com
    '''
    try:
        name = soup.find("div", {"id":"mainReportBox"})\
                   .find("span",{"class":"font18 medium_blue_type w-100 pull-left"})\
                   .text.replace(';','')
    except:
        name = ''
    return name

def etf_description(soup):
    '''
    * soup = make_soup1(ticker)
    * This function will return the description of ETF from ETF.com

    * Decription = etf_description(soup)[0]
    * Analytics  = etf_description(soup)[1]
    '''
    try:
        overview = soup.find("div",{"id":"overview"}).text.split('\n')
        Description = ''.join([s for s in overview if "Fund Description" in s]).split("Fund Description ",)[1]
        Analytics = ''.join([s for s in overview if "Factset Analytics Insight" in s]).split("Factset Analytics Insight ",)[1]
    except:
        Description, Analytics = ''
    return Description, Analytics

def make_fund_summary(soup):
    '''
    * soup = make_soup1(ticker)
    * This function will return the basic summary information of ETF from ETF.com

    * Issuer, Brand, Inception Date, Expense Ratio, Assets Under Management, ...
    * Ex) Issuer = make_fund_summary(soup).loc['Issuer'][0]
    '''
    try:
        summary_info = soup.find("div", {"id":"fundSummaryData"}).get_text()    
        summary_list = list(filter(None, summary_info.split('\n')))
        summary_df = list_to_df(summary_list)
    except:
        summary_df = ''
    return summary_df

def make_portfolio(soup):
    '''
    * soup = make_soup1(ticker)
    * This function will return the portfolio information of ETF from ETF.com

    * for Equity       : 'Distribution Yield', 'Next Ex-Dividend Date', 'Number of Holdings'
    * for Fixed Income : 'Duration', 'Yield to Maturity', 'Next Ex-Dividend Date'
    * for Commodity    : 'Exposure Type', 'Rolling Strategy', 'Rebalancing Frequency'
    * for Alternatives : 'Starategy', 'Distribution Yield', 'Next Ex-Dividend Date'
    '''
    try:
        portfolio_info = soup.find("div", {"id":"fundPortfolioData"}).get_text()    
        portfolio_list = list(filter(None, portfolio_info.split('\n')))
        portfolio_df = list_to_df(portfolio_list)
    except:
        portfolio_df = ''
    return portfolio_df

def make_tradability(soup):
    '''
    * soup = make_soup1(ticker)
    * This function will return the tradability of ETF from ETF.com

    * 'Median Daily Share Volume', 'Average Spread', 
    * 'Median Premium / Discount', 'Max. Premium / Discount'
    '''
    try:
        tradability_info = soup.find("div", {"id":"fundTradabilityData"}).get_text()
        tradability_list = list(filter(None, tradability_info.split('\n')))
        tradability_df = list_to_df(tradability_list)
    except:
        tradability_df = ''
    return tradability_df

def make_index(soup):
    '''
    * soup = make_soup1(ticker)
    * This function will return the benchmark information of ETF from ETF.com

    * 'Index Tracked', 'Index Selection Methodology'
    '''
    try:
        index_info = soup.find("div", {"id":"fundIndexData"}).get_text()
        index_list = list(filter(None, index_info.split('\n')))
        index_df = list_to_df(index_list).loc['Index Tracked'][0]
    except:
        index_df = ''
    return index_df




def top_holdings(soup):
    '''
    * soup = make_soup2(ticker)
    * This function will return top holdings securities of ETF from ETFdb.com
    '''
    try:
        # Note that source from ETFdb.com !- source from ETF.com
        th_tags = soup.find('div', {'id':'holding_section'}).find_all('td')    # find _all
        th_list = [tag.get_text() for tag in th_tags][:-1]
        th_df = make_top_holdings(th_list)
    except:
        th_df = ''
    return th_df



def holding_analysis(soup, num_info):
    '''
    * soup = make_soup3(ticker)
    *
    * num_info
    * For Equity, Real Estate
    * [0]: Asset Allocation
    * [1]: Sector
    * [2]: Market Cap
    * [3]: Region
    * [4]: Market Tier
    * [5]: Country
    *
    * For Fixed Income
    * [0]: Asset
    * [1]: Bond Sector
    * [2]: Maturity
    '''

    try:
        holding_tags = soup.find("div", {"class":"holdings-chart-detailed-section"}).find_all_next('table')[num_info]
        holding_list = list(filter(None, holding_tags.get_text().split('Percentage')[1].split('\n')))
        holding_df = pd.DataFrame(holding_list[1::2], index=holding_list[::2])
    except:
        holding_df = ''

    return holding_df





# %% Plotting Functions

def truncate(df):
    temp = df[df.columns[0]].str.rstrip('%')
    temp_float = pd.DataFrame([float(i) for i in temp])
    temp_float.index = temp.index

    percent_limit = 5

    maindata = temp_float[temp_float >= percent_limit].dropna()

    other_label = "etc. (< " + str(percent_limit) + "%)"
    otherdata = pd.DataFrame(temp_float[temp_float < percent_limit].sum())
    otherdata.index = [other_label]

    data = maindata.append(otherdata)
    return data

def my_autopct(pct):
    return ('%1.1f%%' % pct) if pct > 5 else ''

def plot_pie_chart(df):
    data = truncate(df)   
    labels = data.index.tolist()
    datavals = np.round(data[data.columns[0]].tolist(), 1)
    if len(labels) <= 9:
        cmap = plt.get_cmap("rainbow")
    elif len(labels) > 9:
        cmap = plt.get_cmap("rainbow")
    color = cmap(np.linspace(0., 1., len(labels)))

    fig = plt.figure(figsize=(6.1,4))
    gs = gridspec.GridSpec(1, 1)
    ax = plt.subplot(gs[0])    
    ax.pie(x = datavals, autopct=my_autopct, pctdistance = 0.75,
           startangle = 90, counterclock = False, colors = color,
           wedgeprops=dict(width=0.8, edgecolor='w'), textprops={'fontsize':14, 'color':"w"})
    ax.legend(labels, prop={'size':14}, loc="center left", bbox_to_anchor=(0.9, 0, 0.5, 1),frameon = False)
    plt.tight_layout()
    return fig


def plot_pie_chart2(df):
    data = truncate(df)   
    labels = data.index.tolist()
    datavals = np.round(data[data.columns[0]].tolist(), 1)
    if len(labels) <= 9:
        cmap = plt.get_cmap("rainbow")
    elif len(labels) > 9:
        cmap = plt.get_cmap("rainbow")
    color = cmap(np.linspace(0., 1., len(labels)))

    fig = plt.figure(figsize=(6.1,4))
    gs = gridspec.GridSpec(1, 1)
    ax = plt.subplot(gs[0])    
    ax.pie(x = datavals, autopct=my_autopct, pctdistance = 0.75,
           startangle = 90, counterclock = False, colors = color,
           wedgeprops=dict(width=0.8, edgecolor='w'), textprops={'fontsize':14, 'color':"w"})
    ax.legend(labels, prop={'size':14}, loc="center left", bbox_to_anchor=(0.9, 0, 0.5, 1),frameon = False)
    plt.tight_layout()
    return fig

def plot_hbar_chart(df):
    data = truncate(df)
    labels = data.index.tolist()
    datavals = data[data.columns[0]].tolist()

    labels.reverse()
    datavals.reverse()

    fig, ax = plt.subplots( figsize = (6.2, 4))
    ax.barh(labels, datavals, align = 'center')
    ax.spines['right'].set_visible(False)
    ax.spines['top'].set_visible(False)

    for i, v in enumerate(datavals):
        ax.text(v+0.15, i, str(np.round(v,1)))
    return fig

def month_return_table(df):
    m_ret = pd.DataFrame(df.resample('M').last().pct_change()*100)
    m_ret.columns = ['Returns']
    m_ret_idx = m_ret.index
    m_ret['Year'] = m_ret_idx.strftime('%Y')
    m_ret['Month'] = m_ret_idx.strftime('%b')
    
    m_ret_pivot = m_ret.pivot('Year', 'Month', 'Returns')
    
    # handle missing months
    months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    
    if len(m_ret_pivot.columns) < 12:
        add_months = list(set(months) - set(m_ret_pivot.columns))
        for i in range(len(add_months)):
            m_ret_pivot[add_months[i]] = np.nan
    m_ret_pivot = m_ret_pivot[months]
    
    res = m_ret_pivot[-3:]
    return res

def mret_hmap(df):
    fig, ax = plt.subplots(figsize=(10.7, 3.0))
    sns.set(font_scale = 1.0)
    ax = sns.heatmap(df, ax=ax, annot = True,
                     annot_kws={'size': 14},
                     fmt = '.1f', cmap='Spectral',
                     center=0, vmin = -30, vmax = 30,
                     linewidths = 0.2, linecolor = 'white')
    ax.xaxis.set_ticks_position('bottom')
    ax.tick_params(axis='both', which='both', length=0)
    plt.xlabel('')
    plt.ylabel('')
    plt.yticks(rotation = 0, fontsize=14)
    plt.xticks(rotation = 0, fontsize=14)    
    return fig

def price_plot(ticker):
    df = pdr.get_data_yahoo(ticker, start = start_date, end = end_date)
    Close = df['Close']
    adj_Close = df['Adj Close']
    Adj_Close = adj_Close/adj_Close[0]*Close[0]
    fig = plt.figure(figsize = (9.3, 3.5), frameon=False)
    gs = gridspec.GridSpec(1, 1)    
    ax = plt.subplot(gs[0])
    ax.plot(Adj_Close, color='goldenrod', label = 'Adj Price of %s' %ticker)
    ax.set_facecolor('white')
    ax.legend
    ax.minorticks_off()
    ax.xaxis.set_major_formatter(DateFormatter("%y-%m"))
    plt.xticks(fontsize=15)
    plt.yticks(fontsize=15)

    res = month_return_table(Adj_Close)
    return fig, res



def make_tradibility(ticker):
    df = pdr.get_data_yahoo(ticker, start = start_date2, end = end_date2)
    df = df.rolling(20).mean()
    df = pd.DataFrame(df.resample('M').last())
    Vol = df['Volume']
    Vol = Vol/1000
    fig, ax = plt.subplots(figsize = (6.55, 3.5))
    ax.plot(Vol, color='dodgerblue')
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%y-%m'))
    ax.minorticks_off()
    ax.set_facecolor('white')
    ax.fill_between(Vol.index, Vol, 0, color = 'dodgerblue')
    plt.yticks(rotation = 0, fontsize=15)
    plt.xticks(rotation = 0, fontsize=15)
    plt.tight_layout()

    return fig  



# %% Powerpoint Functions

def text_to_text_placeholder(slide, placeholder_number, txt):
    text_placeholder = slide.shapes[placeholder_number].text_frame
    text_placeholder.text = txt

def top10_df_to_table(slide, df):
    '''
    top 10 holdings 
    '''
    rows, cols = df.shape
    placeholder = slide.shapes[5]   
    res = placeholder.insert_table(rows, cols+1)    
    res.table.columns[0].width = Cm(6.8)
    res.table.columns[1].width = Cm(1.4)
    
    for row in range(rows):
        res.table.rows[row].height = Cm(5.2 / rows)
        
    m = df.index.values.reshape(rows, 1)
    n = df.values
    mat = np.hstack((m, n))

    for row in range(rows):
        for col in range(cols+1):
            val = mat[row, col]
            text = str(val)
            cell = res.table.cell(row, col)
            cell.text = text
            
            cell.text_frame.paragraphs[0].font.size=Pt(8)
            cell.text_frame.paragraphs[0].font.bold = False
            cell.text_frame.paragraphs[0].font.name = '원신한 Light'
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            if col == 0:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            elif col == 1:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            
            cell.fill.background()
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Cm(0.1)
            cell.margin_top  = cell.margin_bottom = Cm(0.0)

def summary_df_to_table(slide, df):
    '''
    the basic information of a etf
    '''
    rows = 9
    cols = 2
    placeholder = slide.shapes[4]
    
    res = placeholder.insert_table(rows, cols)
    
    res.table.columns[0].width = Cm(3.6)
    res.table.columns[1].width = Cm(4.6)
    
    for row in range(rows):
        res.table.rows[row].height = Cm(4.8 / rows)
        
    m = df.index.values.reshape(rows, 1)
    n = df.values

    mat = np.hstack((m, n))

    for row in range(rows):
        
        for col in range(cols):
            val = mat[row, col]
            text = str(val)
            cell = res.table.cell(row, col)
            cell.text = text
            
            cell.text_frame.paragraphs[0].font.size=Pt(8)
            cell.text_frame.paragraphs[0].font.bold = False
            cell.text_frame.paragraphs[0].font.name = '원신한 Light'
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            if col == 0:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            elif col == 1:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            
            cell.fill.background()
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Cm(0.1)
            cell.margin_top  = cell.margin_bottom = Cm(0.0)

def chart_to_picture(slide, placeholder_number, fig_path):
    picture_placeholder = slide.shapes[placeholder_number]
    picture_placeholder.insert_picture(fig_path)






# %% Main Function



Equity = Fixed_Income = Real_Estate = Multi_Asset = Currency = Commodity = Alternative = Sus_Equity = Sus_FI = []


def make_slide(ticker):
    try:
        prs = Presentation(save_path + '/etf_gd_ver1.pptx')    # open a pptx file
    
        if ticker in Equity + Real_Estate + Sus_Equity:
            layout = prs.slide_layouts[1]    # slide layout
        elif ticker in Commodity + Currency + Alternative:
            layout = prs.slide_layouts[2]
        elif ticker in Fixed_Income + Sus_FI:
            layout = prs.slide_layouts[3]
        elif ticker in Multi_Asset:
            layout = prs.slide_layouts[4]
    
        slide = prs.slides.add_slide(layout)
    
        soup1 = make_soup1(ticker)
        soup2 = make_soup2(ticker)
        soup3 = make_soup3(ticker)
           
        # Classification
        # Equity
        if ticker in Equity_WD:
            cl = '주식 > 전세계'
        elif ticker in Equity_DM:
            cl = '주식 > 미국외 선진국'
        elif ticker in Equity_US:
            cl = '주식 > 미국 지수형'
        elif ticker in Equity_US_Smartbeta:
            cl = '주식 > 미국 스마트베타'
        elif ticker in Equity_US_Sec:
            cl = '주식 > 미국 섹터'
        elif ticker in Equity_EM:
            cl = '주식 > 신흥국'
        elif ticker in Equity_Theme:
            cl = '주식 > 테마형'
        #
        # Fixed Income
        elif ticker in FI_Broad:
            cl = '채권 > 종합'
        elif ticker in FI_Govt:
            cl = '채권 > 국채'
        elif ticker in FI_IG:
            cl = '채권 > 투자등급(IG)'
        elif ticker in FI_HY:
            cl = '채권 > 하이일드(HY)'
        elif ticker in FI_MBS:
            cl = '채권 > 주택저당증권(MBS)'
        elif ticker in FI_Inflation:
            cl = '채권 > 물가연동'
        elif ticker in FI_EM:
            cl = '채권 > 신흥국'
        #
        # Commodity
        elif ticker in Comm_Broad:
            cl = '원자재 > 종합'
        elif ticker in Comm_Energy:
            cl = '원자재 > 에너지'
        elif ticker in Comm_Precious:
            cl = '원자재 > 귀금속'
        elif ticker in Comm_Agri:
            cl = '원자재 > 농산물'
        #
        # Real Estate
        elif ticker in RE_Broad:
            cl = '리츠 > 종합'
        elif ticker in RE_Industrial:
            cl = '리츠 > 산업용'
        elif ticker in RE_Data:
            cl = '리츠 > 데이터센터'
        elif ticker in RE_Mortgage:
            cl = '리츠 > 모기지'
        elif ticker in RE_Residential:
            cl = '리츠 > 주거용'
        elif ticker in RE_HC:
            cl = '리츠 > 헬스케어'
        #
        # Alternative
        elif ticker in Alt_HedgeFund:
            cl = '헤지펀드 > 종합'
        elif ticker in Alt_AltIncome:
            cl = '대안인컴'
        #
        # Currency
        elif ticker in Currency:
            cl = '환율'
        #
        # Multi Asset
        elif ticker in Multi_Asset:
            cl = '멀티에셋/자산배분'
        #
        # Sus_Equity
        elif ticker in Sus_Equity:
            cl = '지속가능성 > 주식' 
        #
        # Sus_FI
        elif ticker in Sus_FI:
            cl = '지속가능성 > 채권'              
        else:
            cl = ''
    
      #----------------------------------------------------------------------------
      # upper part of a slide
    
        # 1. Ticker
        text_to_text_placeholder(slide, 0, ticker)
    
        # 2. Name
        text_to_text_placeholder(slide, 1, etf_name(soup1))
        
        # 3. Benchmark
        text_to_text_placeholder(slide, 16, cl + ' > ' + make_index(soup1))
    
        # 4. Analytics
        text_to_text_placeholder(slide, 3, ETF_description.loc[ticker][0])
        
    
    
      #----------------------------------------------------------------------------
      # middle part of a slide
    
        # 1. ETF Details
        df1 = make_fund_summary(soup1)
        df2 = make_portfolio(soup1)
        df3 = make_tradability(soup1)
        summary_df = pd.concat([df1, df2, df3], axis=1)
    
        summary_df_idx = summary_df.index.values
        new_summary_df_idx = [idx.split('(', 1)[0].strip() for idx in summary_df_idx]
        summary_df.index = new_summary_df_idx
        summary_df = summary_df.loc[~summary_df.index.duplicated(keep='first')]
    
        common_idx = ['Issuer', 'Inception Date', 'Assets Under Management', 'Expense Ratio',
                       'Median Daily Share Volume', 'Average Spread',
                       'Median Premium / Discount']
    
        # for Equity and Real Estate ETFs
        equity_idx = ['Distribution Yield', 'Number of Holdings']
    
        # for Fixed Income ETFs
        fixed_income_idx = ['Duration', 'Yield to Maturity']
    
        # for Commodity ETFs
        commodity_idx = ['Exposure Type', 'Rolling Strategy']
    
        # for Alternatives ETFs
        alternative_idx = ['Strategy', 'Distribution Yield']
    
        # for Sus_Equity ETFs
        sus_equity_idx = ['Distribution Yield','Number of Holdings']
        
        # for Sus_FI ETFs
        sus_fi_idx = ['Duration', 'Yield to Maturity']
        
        if ticker in Equity + Real_Estate:
            summary_idx = common_idx + equity_idx
        elif ticker in Fixed_Income:
            summary_idx = common_idx + fixed_income_idx
        elif ticker in Commodity:
            summary_idx = common_idx + commodity_idx
        elif ticker in Alternative:
            summary_idx = common_idx + alternative_idx
        elif ticker in Sus_Equity: 
            summary_idx = common_idx + sus_equity_idx
        elif ticker in Sus_FI: 
            summary_idx = common_idx + sus_fi_idx
        else:
            summary_idx = common_idx + ['', '']
    
        val = []
        for i in summary_idx:
            try:
                val.append(summary_df.loc[i].dropna()[0].strip())
            except:
                val.append('')
    
        details_df = pd.DataFrame(val, index=summary_idx)
        common_idx_ = ['운용사', '상장일', '자산규모', '총비용',
                       '일평균 거래량', '평균 호가 스프레드',
                       'NAV 대비 스프레드']
        
        # for Equity and Real Estate ETFs
        equity_idx_ = ['분배율', '편입종목수']
    
        # for Fixed Income ETFs
        fixed_income_idx_ = ['듀레이션', 'YTM']
    
        # for Commodity ETFs
        commodity_idx_ = ['투자방법', '롤오버 전략']
    
        # for Alternatives ETFs
        alternative_idx_ = ['투자전략', '분배율']
    
        # for Sus_Equity ETFs
        sus_equity_idx_ = ['분배율', '편입종목수']
    
        # for Sus_FI ETFs
        sus_fi_idx_ = ['듀레이션', 'YTM']
    
    
        if ticker in Equity + Real_Estate:
            summary_idx_ = common_idx_ + equity_idx_
        elif ticker in Fixed_Income:
            summary_idx_ = common_idx_ + fixed_income_idx_
        elif ticker in Commodity:
            summary_idx_ = common_idx_ + commodity_idx_
        elif ticker in Alternative:
            summary_idx_ = common_idx_ + alternative_idx_
        elif ticker in Sus_Equity:
            summary_idx_ = common_idx_ + sus_equity_idx_
        elif ticker in Sus_FI:
            summary_idx_ = common_idx_ + sus_fi_idx_      
        else:
            summary_idx_ = common_idx_ + ['', '']
        
        details_df.index = summary_idx_
        
        summary_df_to_table(slide, details_df)
    
    
        # 2. Top 10 Holdings
        top10_df = top_holdings(soup2)[:10]
        top10_df_to_table(slide, top10_df)
    
    
        # 3. Historical Price and Monthly Returns Table
        a, b = price_plot(ticker)
    
        hp_fig = a
        hp_fig_path = save_path + '/hp_fig.png'
        hp_fig.savefig(hp_fig_path, bbox_inches = 'tight')
        chart_to_picture(slide, 6, hp_fig_path)
    
        mr_fig = mret_hmap(b)
        mr_fig_path = save_path + '/mr_fig.png'
        mr_fig.savefig(mr_fig_path, bbox_inches = 'tight')
        chart_to_picture(slide, 15, mr_fig_path)
    
      #----------------------------------------------------------------------------
      # Bottop part of a slide
    
        if ticker in Equity + Real_Estate + Sus_Equity:
            # Sector
            text_to_text_placeholder(slide, 11, '업종 분포')
            sc_df = holding_analysis(soup3, 2)
            sc_fig = plot_pie_chart(sc_df)
            sc_fig_path = save_path + '/sc_fig.png'
            sc_fig.savefig(sc_fig_path, bbox_inches='tight')
            chart_to_picture(slide, 12,  sc_fig_path)
            
            # Trading Volume   
            text_to_text_placeholder(slide, 7, '월간 평균 거래량(천 주)')
            v = make_tradibility(ticker)
            vp_fig = v
            vp_fig_path = save_path + '/vp_fig.png'
            vp_fig.savefig(vp_fig_path, bbox_inches = 'tight')
            chart_to_picture(slide, 8, vp_fig_path)    

        
        elif ticker in Fixed_Income + Sus_FI:
            # Asset Allocation
            text_to_text_placeholder(slide, 11, '자산 구성')
            cp_df = holding_analysis(soup3, 0)
            cp_fig = plot_pie_chart(cp_df)
            cp_fig_path = save_path + '/cp_fig.png'
            cp_fig.savefig(cp_fig_path, bbox_inches='tight')
            chart_to_picture(slide, 12, cp_fig_path)
            # Maturity
            #text_to_text_placeholder(slide, 9, '만기 분포')
            #cp_df = holding_analysis(soup3, 2)
            #cp_fig = plot_pie_chart(cp_df)
            #cp_fig_path = save_path + '/cp_fig.png'
            #cp_fig.savefig(cp_fig_path, bbox_inches='tight')
            #chart_to_picture(slide, 10, cp_fig_path)     

            # Trading Volume   
            text_to_text_placeholder(slide, 7, '월간 평균 거래량(천 주)')
            v = make_tradibility(ticker)
            vp_fig = v
            vp_fig_path = save_path + '/vp_fig.png'
            vp_fig.savefig(vp_fig_path, bbox_inches = 'tight')
            chart_to_picture(slide, 8, vp_fig_path)                      
            
        elif ticker in Multi_Asset:
            # Asset Allocation
            text_to_text_placeholder(slide, 13, '자산배분 현황')
            aa_df = holding_analysis(soup3, 0)
            aa_fig = plot_pie_chart2(aa_df)
            aa_fig_path = save_path + '/aa_fig.png'
            aa_fig.savefig(aa_fig_path, bbox_inches='tight')
            chart_to_picture(slide, 14,  aa_fig_path)
            # Equity Sector
            text_to_text_placeholder(slide, 11, '시가총액 구분')
            sc_df = holding_analysis(soup3, 2)
            sc_fig = plot_pie_chart(sc_df)
            sc_fig_path = save_path + '/sc_fig.png'
            sc_fig.savefig(sc_fig_path, bbox_inches='tight')
            chart_to_picture(slide, 12, sc_fig_path)
            
            # Trading Volume   
            text_to_text_placeholder(slide, 7, '월간 평균 거래량(천 주)')
            v = make_tradibility(ticker)
            vp_fig = v
            vp_fig_path = save_path + '/vp_fig.png'
            vp_fig.savefig(vp_fig_path, bbox_inches = 'tight')
            chart_to_picture(slide, 8, vp_fig_path)   
            
        elif ticker in Commodity + Alternative + Currency:
            # Trading Volume   
            text_to_text_placeholder(slide, 7, '월간 평균 거래량(천 주)')
            v = make_tradibility(ticker)
            vp_fig = v
            vp_fig_path = save_path + '/vp_fig.png'
            vp_fig.savefig(vp_fig_path, bbox_inches = 'tight')
            chart_to_picture(slide, 8, vp_fig_path)  
            
            
            
        
            
        else:
            pass
        
        
 
        
        prs.save(save_path + '/etf_gd_ver1.pptx')
        print("%s: complete" %ticker)
        plt.clf()
    except:
        print('%s: Fail' %ticker)



# %% Classification
Equity_US = []
Equity_US_Smartbeta = []
Equity_US_Sec = []
Equity_WD = []
Equity_DM = []
Equity_EM = []
Equity_Theme = ['FXD', 'IYT', 'KRE', 'IQLT', 'MOAT', 'MGK', 'MILN', 'XT', 'SPYD', 'VIG', 'SOXX', 'WDIV', 'LIT', 'PIE']

Equity = Equity_US + Equity_US_Smartbeta + Equity_US_Sec + Equity_WD +\
         Equity_DM + Equity_EM + Equity_Theme

FI_Broad = []
FI_Govt = []
FI_IG = ['VCSH']
FI_HY = ['HYD','SHYG', 'BKLN', 'SRLN']
FI_MBS = []
FI_Inflation = ['TIP']
FI_EM = []
Fixed_Income = FI_Broad + FI_Govt + FI_IG + FI_HY + FI_MBS + FI_Inflation + FI_EM

Comm_Broad = []
Comm_Energy = []
Comm_Precious = []
Comm_Agri = []
Commodity = Comm_Broad + Comm_Energy + Comm_Precious + Comm_Agri

RE_Broad = ['VNQ']
RE_Industrial = []
RE_Data = []
RE_Mortgage = []
RE_Residential = []
RE_HC = []
Real_Estate = RE_Broad + RE_Industrial + RE_Data + RE_Mortgage + RE_Residential + RE_HC

Alt_HedgeFund = []
Alt_AltIncome = ['CWB', 'PFF', 'IGF', 'BIZD', 'PCEF']

Alternative = Alt_HedgeFund + Alt_AltIncome

Currency = []

Multi_Asset = []

Sus_Equity = []

Sus_FI = []


univ = Equity + Fixed_Income + Commodity + Real_Estate + Alternative + Currency + Multi_Asset + Sus_Equity + Sus_FI



# %% Main

for i in Equity_US:
    ticker = i
    make_slide(ticker)

for i in Equity_US_Smartbeta:
    ticker = i
    make_slide(ticker)
    
for i in Equity_US_Sec:
    ticker = i
    make_slide(ticker)
    
for i in Equity_WD:
    ticker = i
    make_slide(ticker)
    
for i in Equity_DM:
    ticker = i
    make_slide(ticker)
    
for i in Equity_EM:
    ticker = i
    make_slide(ticker)
    
for i in Equity_Theme:
    ticker = i
    make_slide(ticker)

for i in Fixed_Income:
    ticker = i
    make_slide(ticker)

for i in Commodity:
    ticker = i
    make_slide(ticker)

for i in Real_Estate:    
    ticker = i
    make_slide(ticker)

for i in Alternative:
    ticker = i
    make_slide(ticker)

for i in Currency:
    ticker = i
    make_slide(ticker)

for i in Multi_Asset:
    ticker = i
    make_slide(ticker)

for i in Sus_Equity:
    ticker = i
    make_slide(ticker)
    
for i in Sus_FI:
    ticker = i
    make_slide(ticker)
    


