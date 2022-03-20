# -*- coding: utf-8 -*-
"""
Created on Wed Feb  3 15:51:06 2021
ㅇ
@author: Seunghoon Baek
"""


from bs4 import BeautifulSoup
import requests
import pandas as pd
import numpy as np
import re
from pandas_datareader import data as pdr
import yfinance as yf
yf.pdr_override()
from datetime import datetime, timedelta
import datetime as dt
import matplotlib.pyplot as plt 
import matplotlib.ticker as tick
import matplotlib.gridspec as gridspec
import matplotlib.dates as mdates
from matplotlib.dates import DateFormatter
from matplotlib.dates import MO
plt.rcParams.update({'font.size':15})
plt.rcParams.update({'figure.max_open_warning': 0})
plt.rcParams['axes.unicode_minus'] = False
import math
import seaborn as sns
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

import random
from selenium import webdriver


# %% Data From Excel
save_path = r'C:/ETF'
df_from_excel = pd.read_excel(save_path + '\KD ETF Description.xlsx',
                              index_col = 'symbol')
ETF_description = df_from_excel.copy().dropna()



# %% save path

save_path = r'C:/ETF'

# %% Date Setting

def make_beg_end_date(now):
    today = dt.date(now.year, now.month, now.day)    
    start_date = '1991-01-01'

    if today.strftime('%A') == 'Saturday':
        end_date = today - timedelta(days=1)    
    elif today.strftime('%A') == 'Sunday':
        end_date = today - timedelta(days=2)    
    else:
        end_date = today

    return start_date, end_date

now = datetime.now()
start_date, end_date = make_beg_end_date(now)

# %% Crawling Functions

ticker = 'ticker'
NonBreakSpace = u'\xa0'
years = mdates.YearLocator()
months = mdates.MonthLocator()
weeks = mdates.WeekdayLocator(byweekday=MO, interval=4)

def tag_to_list(tag):
    temp = tag.get_text().split('\n')
    return_list = list(filter(None, temp))
    return return_list

def list_to_df(input_list):
    df_columns = re.sub(r'\[[^)]*\]', '', input_list[0])
    df_index = input_list[1::2]
    df = pd.DataFrame(input_list[2::2], index = df_index, columns = [df_columns])
    df.index = df.index.str.replace(';','')
    return df

def make_soup1(ticker):
    source1 = "https://navercomp.wisereport.co.kr/v2/ETF/Index.aspx?cn=&menuType=block&cmp_cd="
    url1 = source1 + ticker[1:]
    driver = webdriver.Chrome('C:\ETF\chromedriver.exe')
    driver.implicitly_wait(random.uniform(7, 10))
    driver.get(url1)
    html = driver.page_source
    driver.quit()
    soup1 = BeautifulSoup(html, 'html.parser')
    return soup1

def make_soup2(ticker):
    source2 = "http://comp.fnguide.com/SVO2/ASP/etf_snapshot.asp?pGB=1&cID=&MenuYn=Y&ReportGB=&NewMenuID=401&stkGb=770&gicode=A"
    url2 = source2 + ticker[1:]
    driver = webdriver.Chrome('C:\ETF\chromedriver.exe')
    driver.implicitly_wait(random.uniform(7, 10))
    driver.get(url2)
    html = driver.page_source
    driver.quit()
    soup2 = BeautifulSoup(html, 'html.parser')
    return soup2

def etf_name(soup):
    '''
    * soup = make_soup2(ticker)
    * This function will return the full name of ETF from fnguide
    '''
    try:
        name = soup.find("div", {"class":"contwrap"})\
                   .find("h1")\
                   .text.replace(NonBreakSpace,' ')
    except:
        name = ''
    return name

def make_index(soup):
    '''
    * soup = make_soup2(ticker)
    * This function will return the benchmark information of ETF from fnguide

    '''
    try:
        index_df = soup.find("div", {"class":"contwrap"})\
                       .find("dd",{"class":"fl bench"})\
                       .text.replace(NonBreakSpace, ' ')
    except:
        index_df = ''
    return index_df

def etf_description(soup):
    '''
    * soup = make_soup1(ticker)
    * This function will return the description of ETF from naver finance

    '''
    try:
        description = soup.find("div", {"class":"cmp_comment"}).text.replace('\n', '').replace('\t', '')
        
    except:
        description = ''
    return description

def top_holdings(soup):
    '''
    * soup = make_soup2(ticker)
    * This function will return top holdings securities of ETF from fnguide
    '''
    try:
        th_tags1 = soup.find('div', {'id':'gridTop10'}).find_all('th')    # holdings
        th_tags2 = soup.find('div', {'id':'gridTop10'}).find_all('td')    # allocations
        th_list1 = [tag.get_text() for tag in th_tags1][2:]
        th_list2 = [tag.get_text() for tag in th_tags2]
        th_df = pd.DataFrame(th_list2, th_list1)
    except:
        th_df = ''
    return th_df


def make_fund_summary1(soup):
    '''
    * soup = make_soup1(ticker)
    * This function will return the fund summary of ETF from naver finance

    '''
    try:
        mf_tags1 = soup.find('tbody', {'id':'status_grid_body'}).find_all('th')
        mf_tags2 = soup.find('tbody', {'id':'status_grid_body'}).find_all('td')
        mf_list1 = [tag.get_text() for tag in mf_tags1]
        mf_list2 = [tag.get_text() for tag in mf_tags2]
        mf_comment1 = [comment.replace(' ', '') for comment in mf_list1]
        mf_comment2 = [comment.replace(' ', '') for comment in mf_list2]
        df1 = pd.DataFrame(mf_comment2, mf_comment1)

    except:
        df1 = ''
    return df1
    
def make_fund_summary2(soup):
    '''
    * soup = make_soup1(ticker)
    * This function will return the fund summary of ETF from naver finance

    '''
    try:
        mf_tags3 = soup.find('tbody', {'id':'product_summary_grid_body'}).find_all('th')
        mf_tags4 = soup.find('tbody', {'id':'product_summary_grid_body'}).find_all('td')
        mf_list3 = [tag.get_text() for tag in mf_tags3]
        mf_list4 = [tag.get_text() for tag in mf_tags4]
        mf_comment3 = [comment.replace(' ', '') for comment in mf_list3]
        mf_comment4 = [comment.replace(' ', '') for comment in mf_list4]        
        df2 = pd.DataFrame(mf_comment4, mf_comment3)

    except:
        df2 = ''
    return df2  

def make_fund_summary3(soup):
    '''
    * soup = make_soup1(ticker)
    * This function will return the fund summary of ETF from naver finance

    '''
    try:
        mf_tags5 = soup.find('tbody', {'id':'NAV_grid_body'}).find('tr').find_all('td')
        mf_list5 = [tag.get_text() for tag in mf_tags5]
        mf_list5 = mf_list5[1:2]
        df2_1 = pd.DataFrame(mf_list5, index = ['최근 종가(원)'])

    except:
        df2_1 = ''
    return df2_1     


def make_tradibility(soup):
    '''
    * soup = make_soup1(ticker)
    * This function will return the tradibility of ETF from naver finance,

    '''
    try:
        mt_tags1 = soup.find('tbody', {'id':'volume_grid_body'})\
                       .find_all('th')      
        mt_tags2 = soup.find('tbody', {'id':'volume_grid_body'})\
                       .find_all('td')
        mt_list1 = [tag.get_text() for tag in mt_tags1]
        mt_list2 = [tag.get_text() for tag in mt_tags2]
        
        if len(mt_list1)>=12:
            mt_list2_1 = mt_list2[0::2]
            mt_list2_1.reverse()
            mt_list1.reverse()
            mt_list1 = pd.to_datetime(mt_list1)
            mt_comment2_1 = [comment.replace('-', '0') for comment in mt_list2_1]
            mt_comment2_2 = [comment.replace(',', '') for comment in mt_comment2_1]
            df3 = pd.DataFrame(mt_comment2_2, mt_list1, columns=['거래량(천주)'])
            df3['거래량(천주)'] = pd.to_numeric(df3['거래량(천주)'])
            fig, ax = plt.subplots(figsize=(6.4, 3.5))
            ax.plot(df3['거래량(천주)'], color='dodgerblue')
            ax.legend(frameon = True)
            ax.xaxis.set_major_locator(mdates.MonthLocator(interval=3))
            ax.xaxis.set_major_formatter(mdates.DateFormatter('%y-%m'))
            ax.minorticks_off()
            ax.set_facecolor('white')        
            ax.fill_between(df3['거래량(천주)'].index, df3['거래량(천주)'], 0, color='dodgerblue')
            plt.yticks(rotation = 0, fontsize=15)
            plt.xticks(rotation = 0, fontsize=15)    
#            plt.grid(linestyle="--", color='lightgrey')
            plt.tight_layout()
            
        else:
            mt_list2 = list(dict.fromkeys(mt_list2))
            mt_list2_1 = mt_list2[0::2]

            try:
                mt_list2_1.remove('\xa0')
                
            except:
                pass
            

            mt_list1 = list(dict.fromkeys(mt_list1))


            try:
                mt_list1.remove('\xa0')
            
            except:
                pass
            
            mt_list1.reverse()
            mt_list1 = pd.to_datetime(mt_list1)
            mt_list2_1.reverse()
            mt_comment2_1 = [comment.replace('-', '0') for comment in mt_list2_1]
            mt_comment2_2 = [comment.replace(',', '') for comment in mt_comment2_1]
            df3 = pd.DataFrame(mt_comment2_2, mt_list1, columns=['거래량(천주)'])
            df3['거래량(천주)'] = pd.to_numeric(df3['거래량(천주)'])
            fig, ax = plt.subplots(figsize=(6.4, 3.5))
            ax.plot(df3['거래량(천주)'], color='cornflowerblue')
            ax.legend(frameon = True)
            ax.xaxis.set_major_locator(mdates.MonthLocator(interval=3))
            ax.xaxis.set_major_formatter(mdates.DateFormatter('%y-%m'))
            ax.minorticks_off()
            ax.set_facecolor('white')        
            ax.fill_between(df3['거래량(천주)'].index, df3['거래량(천주)'], 0, color='cornflowerblue')
            plt.yticks(rotation = 0, fontsize=15)
            plt.xticks(rotation = 0, fontsize=15)    
#            plt.grid(linestyle="--", color='lightgrey')
            plt.tight_layout()
        
    except:
        df3 = ''
    return fig  

def make_tracking_error(soup):
    '''
    * soup = make_soup1(ticker)
    * This function will return the traking error of ETF from naver finance

    '''
    try:
        mte_tags1 = soup.find('tbody', {'id':'NAV_grid_body'})\
                       .find_all('th')      
        mte_tags2 = soup.find('tbody', {'id':'NAV_grid_body'})\
                       .find_all('td')
        mte_list1 = [tag.get_text() for tag in mte_tags1]
        mte_list2 = [tag.get_text() for tag in mte_tags2]
        mte_list2_1 = mte_list2[0::3]
        mte_list2_2 = mte_list2[1::3]
        mte_list1.reverse()
        mte_list1 = pd.to_datetime(mte_list1)
        mte_list2_1.reverse()
        mte_list2_2.reverse()
        mte_comment2_1 = [comment.replace(',', '') for comment in mte_list2_1]
        mte_comment2_2 = [comment.replace(',', '') for comment in mte_list2_2]
        data = list(zip(mte_comment2_1, mte_comment2_2))
        df4 = pd.DataFrame(data, mte_list1, columns = ['순자산가치(NAV)', 'ETF종가'])
        df4['순자산가치(NAV)'] = pd.to_numeric(df4['순자산가치(NAV)'])
        df4['ETF종가'] = pd.to_numeric(df4['ETF종가'])
        df4_1 = 100 * (df4['ETF종가']-df4['순자산가치(NAV)'])/df4['순자산가치(NAV)']
        fig, ax = plt.subplots(figsize=(6.2, 4))
        ax.bar(df4_1.index, df4_1, label='Tracking Error(%)', color='dodgerblue')
        ax.xaxis.set_major_locator(months)
        ax.xaxis.set_major_formatter(mdates.DateFormatter('%y-%m'))
        ax.set_facecolor('white')
        plt.yticks(rotation = 0, fontsize=14)
        plt.xticks(rotation = 0, fontsize=14)    
        plt.tight_layout()
              
    except:
        df4_1 = ''
    return fig  

def holding_analysis(soup):
    '''
    * soup = make_soup2(ticker)

    '''
    try:
        #h_tags1 = soup.find("div", {"class":"um_table", "id":"etf1StockInfo"})\
        #                   .find("tbody").find_all("th")
        h_tags2 = soup.find("div", {"class":"um_table", "id":"etf1StockInfo"})\
                           .find("tbody").find_all("td")
        #h_list1 = [tag.get_text() for tag in h_tags1]
        h_list2 = [tag.get_text() for tag in h_tags2]
        h_list2_1 = h_list2[0::4]
        h_comment2 = [comment.replace(NonBreakSpace, '0') for comment in h_list2_1]
        df5 = pd.DataFrame(h_comment2, index = ['Energy', 'Materials', 'Industrials', 'Consumer Cyclical',\
                                                'Consumer Staples', 'Health Care', 'Financials', 'IT', 'Communications', 'Utilities', 'etc'],\
                           columns = ['펀드'])
        df5['펀드'] = pd.to_numeric(df5['펀드'])
        
        labels = df5['펀드'].index.tolist()

        if len(labels) <= 9:
            cmap = plt.get_cmap("rainbow")
        elif len(labels) > 9:
            cmap = plt.get_cmap("rainbow")       
            
        color = cmap(np.linspace(0., 1., len(labels)))   
        fig = plt.figure(figsize=(6.1,4))
        gs = gridspec.GridSpec(1, 1)
        ax = plt.subplot(gs[0])    
        ax.pie(x = df5['펀드'], autopct=my_autopct, pctdistance = 0.75,
                   startangle = 90, counterclock = False, colors = color,
                   wedgeprops=dict(width=0.8, edgecolor='w'), textprops={'fontsize':14, 'color':"w"})
        ax.legend(labels, prop={'size':14}, loc="center left", bbox_to_anchor=(0.9, 0, 0.5, 1),frameon = False)
        plt.tight_layout()

    except:
        df5 = ''

    return fig


def multiples(soup):
    '''
    * soup = make_soup2(ticker)
    * This function will return multiples of ETF from fnguide
    '''
    try:
        m_tags1 = soup.find("div", {"class":"um_table", "id":"etf1StyleInfoStk"})\
                           .find("tbody").find_all("th")
        m_tags2 = soup.find("div", {"class":"um_table", "id":"etf1StyleInfoStk"})\
                           .find("tbody").find_all("td")                                              
        m_list1 = [tag.get_text() for tag in m_tags1]
        m_comment1 = [comment.replace('대형', 'Large Cap') for comment in m_list1]
        m_comment1_1 = [comment.replace('중소형', 'Mid/Small Cap') for comment in m_comment1]
        m_list2 = [tag.get_text() for tag in m_tags2][0::2]
        m_list3 = [tag.get_text() for tag in m_tags2][1::2]
        m_list2_1 = [comment.replace(NonBreakSpace, '0.00') for comment in m_list2]
        m_list3_1 = [comment.replace(NonBreakSpace, '0.00') for comment in m_list3]
        data = list(zip(m_list2_1, m_list3_1))
        df6 = pd.DataFrame(data, m_comment1_1, columns=['Fund', 'Peers'])
        fig, ax = plt.subplots(figsize=(6.2, 4))
        fig.patch.set_visible(False)
        ax.axis('off')
        y= ax.table(cellText=df6.values, colLabels=df6.columns, rowLabels=df6.index,\
                 loc='center', rowColours=['dodgerblue']*4, colColours=['dodgerblue']*4,\
                 )
        y.auto_set_font_size(False)
        y.set_fontsize(18)
        y.scale(1,5)
        fig.tight_layout()
        

    except:
        df6 = ''
    return fig

def bond_holding_analysis(soup):
    '''
    * soup = make_soup2(ticker)

    '''
    try:
        bh_tags1 = soup.find("div", {"class":"um_table", "id":"etf1BondInfo"})\
                           .find("tbody").find_all("td")
        bh_list1 = [tag.get_text() for tag in bh_tags1]
        bh_list1_1 = bh_list1[0::4]
        bh_list1_2 = bh_list1_1[0:10]
        bh_comment1 = [comment.replace(NonBreakSpace, '0') for comment in bh_list1_2]
        df7 = pd.DataFrame(bh_comment1, index = ['Government', 'Municipal', 'Public Finance', 'Monetary Stabilization',\
                                                'Bank', 'Other Finance', 'Corporate', 'ABS', 'Miscellaneous', 'RF'],\
                           columns = ['편입'])
        df7['편입'] = pd.to_numeric(df7['편입'])
        
        labels = df7['편입'].index.tolist()

        if len(labels) <= 9:
            cmap = plt.get_cmap("rainbow")
        elif len(labels) > 9:
            cmap = plt.get_cmap("rainbow")       

        color = cmap(np.linspace(0., 1., len(labels)))               
#        color = cmap(np.arange(len(labels)))
        fig = plt.figure(figsize=(6.2,4))
        gs = gridspec.GridSpec(1, 1)
        ax = plt.subplot(gs[0])    
        ax.pie(x = df7['편입'], autopct=my_autopct, pctdistance = 0.75,
                   startangle = 90, counterclock = False, colors = color,
                   wedgeprops=dict(width=0.8, edgecolor='w'), textprops={'fontsize':14, 'color':"w"})
        ax.legend(labels, prop={'size':14}, loc="center left", bbox_to_anchor=(1, 0, 0.5, 1),frameon = False)
        
        plt.tight_layout()

    except:
        df7 = ''

    return fig

def credit_analysis(soup):
    '''
    * soup = make_soup2(ticker)
    * This function will return multiples of ETF from fnguide
    '''
    try:
        c_tags1 = soup.find("div", {"class":"um_table", "id":"etf1StyleInfoBond"})\
                           .find("tbody").find_all("td")                                              
        c_list1 = [tag.get_text() for tag in c_tags1]
        c_list1_1 = c_list1[0::2]
        c_list1_2 = c_list1[1::2]
        c_list1_1_1 = c_list1_1[0:3]
        c_list1_2_1 = c_list1_2[0:3]
        data = list(zip(c_list1_1_1, c_list1_2_1))
        df8 = pd.DataFrame(data, index=['Credit Rating', 'Duration', 'Maturity'], columns=['Fund', 'Peers'])
        fig, ax = plt.subplots(figsize=(6.2, 4))
        fig.patch.set_visible(False)
        ax.axis('off')
        z= ax.table(cellText=df8.values, colLabels=df8.columns, rowLabels=df8.index,\
                 loc='center', rowColours=['dodgerblue']*4, colColours=['dodgerblue']*4,\
                 )
        z.auto_set_font_size(False)
        z.set_fontsize(18)
        z.scale(1,4)
        
        fig.tight_layout()
        
    except:
        df8 = ''
    return fig



# %% powerpoint functions

def text_to_text_placeholder(slide, placeholder_number, txt):
    text_placeholder = slide.shapes[placeholder_number].text_frame
    text_placeholder.text = txt

    
def chart_to_picture(slide, placeholder_number, fig_path):
    picture_placeholder = slide.shapes[placeholder_number]
    picture_placeholder.insert_picture(fig_path)


def top10_df_to_table(slide, df):
    '''
    top 10 holdings 
    '''
    rows, cols = df.shape
    placeholder = slide.shapes[5]   
    res = placeholder.insert_table(rows, cols+1)    
    res.table.columns[0].width = Cm(6.8)
    res.table.columns[1].width = Cm(1.4)
    
    for row in range(rows):
        res.table.rows[row].height = Cm(5.2 / rows)
        
    m = df.index.values.reshape(rows, 1)
    n = df.values
    mat = np.hstack((m, n))

    for row in range(rows):
        for col in range(cols+1):
            val = mat[row, col]
            text = str(val)
            cell = res.table.cell(row, col)
            cell.text = text
            
            cell.text_frame.paragraphs[0].font.size=Pt(8)
            cell.text_frame.paragraphs[0].font.bold = False
            cell.text_frame.paragraphs[0].font.name = '원신한 Light'
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            if col == 0:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            elif col == 1:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            
            cell.fill.background()
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Cm(0.1)
            cell.margin_top  = cell.margin_bottom = Cm(0.0)

    
def summary_df_to_table(slide, df):
    '''
    the basic information of a etf
    '''
    rows = 9
    cols = 2
    placeholder = slide.shapes[4]
    
    res = placeholder.insert_table(rows, cols)
    
    res.table.columns[0].width = Cm(3.6)
    res.table.columns[1].width = Cm(4.6)
    
    for row in range(rows):
        res.table.rows[row].height = Cm(4.82 / rows)
        
    m = df.index.values.reshape(rows, 1)
    n = df.values

    mat = np.hstack((m, n))

    for row in range(rows):
        
        for col in range(cols):
            val = mat[row, col]
            text = str(val)
            cell = res.table.cell(row, col)
            cell.text = text
            
            cell.text_frame.paragraphs[0].font.size=Pt(8)
            cell.text_frame.paragraphs[0].font.bold = False
            cell.text_frame.paragraphs[0].font.name = '원신한 Light'
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            
            if col == 0:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            elif col == 1:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            
            cell.fill.background()
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Cm(0.1)
            cell.margin_top  = cell.margin_bottom = Cm(0.0)
    

# %% Plotting Functions

def truncate(df):
    temp = df[df.columns[0]].str.rstrip('%')
    temp_float = pd.DataFrame([float(i) for i in temp])
    temp_float.index = temp.index

    percent_limit = 5

    maindata = temp_float[temp_float >= percent_limit].dropna()

    other_label = "etc. (< " + str(percent_limit) + "%)"
    otherdata = pd.DataFrame(temp_float[temp_float < percent_limit].sum())
    otherdata.index = [other_label]

    data = maindata.append(otherdata)
    return data

def my_autopct(pct):
    return ('%1.1f%%' % pct) if pct > 5 else ''


def month_return_table(df):
    m_ret = pd.DataFrame(df.resample('M').last().pct_change()*100)
    m_ret.columns = ['Returns']
    m_ret_idx = m_ret.index
    m_ret['Year'] = m_ret_idx.strftime('%Y')
    m_ret['Month'] = m_ret_idx.strftime('%b')
    
    m_ret_pivot = m_ret.pivot('Year', 'Month', 'Returns')
    
    # handle missing months
    months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    
    if len(m_ret_pivot.columns) < 12:
        add_months = list(set(months) - set(m_ret_pivot.columns))
        for i in range(len(add_months)):
            m_ret_pivot[add_months[i]] = np.nan
    m_ret_pivot = m_ret_pivot[months]
    
    res = m_ret_pivot[-3:]
    return res

def mret_hmap(df):
    fig, ax = plt.subplots(figsize=(10.7, 3.0))
    sns.set(font_scale = 1.0)
    ax = sns.heatmap(df, ax=ax, annot = True,
                     annot_kws={'size': 14},
                     fmt = '.1f', cmap='Spectral',
                     center=0, vmin = -30, vmax = 30,
                     linewidths = 0.2, linecolor = 'white')
    ax.xaxis.set_ticks_position('bottom')
    ax.tick_params(axis='both', which='both', length=0)
    plt.xlabel('')
    plt.ylabel('')
    plt.yticks(rotation = 0, fontsize=14)
    plt.xticks(rotation = 0, fontsize=14)    
    return fig

def get_price(ticker):
    # count = trading days
    url = "https://fchart.stock.naver.com/sise.nhn?symbol={}&timeframe=day&count=1000&requestType=0".format(ticker)
    get_result = requests.get(url)
    bs_obj = BeautifulSoup(get_result.content, "html.parser")
    
    # information
    inf = bs_obj.select('item')
    columns = ['Date', 'Open' ,'High', 'Low', 'Close', 'Volume']
    df_inf = pd.DataFrame([], columns = columns, index = range(len(inf)))
    
    for i in range(len(inf)):
        df_inf.iloc[i] = str(inf[i]['data']).split('|')
    
    df_inf.index = pd.to_datetime(df_inf['Date'])
    df_inf = df_inf[:-1]
    return df_inf.drop('Date', axis=1).astype(float)

def price_plot2(ticker):
    df = get_price(ticker)
    Close = df['Close']
    fig = plt.figure(figsize = (8.8, 3.5), frameon=False)
    gs = gridspec.GridSpec(1, 1)
    ax = plt.subplot(gs[0])
    ax.plot(Close, color='darkgoldenrod', label = 'Price')
    ax.set_facecolor('white')
    ax.minorticks_off()
    ax.xaxis.set_major_formatter(DateFormatter("%y-%m"))
    plt.xticks(fontsize=15)
    plt.yticks(fontsize=15)
    
    res = month_return_table(Close)
    return fig, res
    

    

# %% make slides

Equity = Equity_Global = Fixed_Income = Real_Estate = Multi_Asset = Currency = Commodity = New = []

def make_slide(ticker):
    try:
        prs = Presentation(save_path + '/etf_kd_ver1.pptx')    # open a pptx file

        if ticker in Equity + Real_Estate + Multi_Asset + Equity_Global + New:
            layout = prs.slide_layouts[1]    # slide layout
        elif ticker in Commodity + Currency:
            layout = prs.slide_layouts[2]
        elif ticker in Fixed_Income:
            layout = prs.slide_layouts[3]

        
        slide = prs.slides.add_slide(layout)
        
        soup1 = make_soup1(ticker)
        soup2 = make_soup2(ticker)
    
      #----------------------------------------------------------------------------
      # upper part of a slide
    
        # 1. Ticker
        text_to_text_placeholder(slide, 0, ticker)
    
        # 2. Name
        text_to_text_placeholder(slide, 1, etf_name(soup2))
    
        # 3. Benchmark
#        text_to_text_placeholder(slide, 2, make_index(soup2))
    
        # 4. Analytics
        text_to_text_placeholder(slide, 3, ETF_description.loc[ticker][0])
    
            
      #----------------------------------------------------------------------------
      # middle part of a slide
      
        # 1. ETF Details
        df1 = make_fund_summary1(soup1)
        df2 = make_fund_summary2(soup1)
        df2_1 = make_fund_summary3(soup1)
        summary_df = pd.concat([df2_1, df1, df2])
        summary_df_new = summary_df.iloc[[0, 2, 5, 6, 8, 9, 11, 13, 16], :]        
        summary_df_to_table(slide, summary_df_new)
        
        
        # 2. Top 10 Holdings
        top10_df = top_holdings(soup2)
        top10_df_to_table(slide, top10_df)
                    

        
        # 3. Historical Price and Monthly Returns Table
        if ticker in Equity + Equity_Global + Real_Estate + Multi_Asset + Commodity + Currency + Fixed_Income + New: 
            tickers = ticker[1:]
            a, b = price_plot2(tickers)
    
            hp_fig = a
            hp_fig_path = save_path + '/hp_fig.png'
            hp_fig.savefig(hp_fig_path, bbox_inches = 'tight')
            chart_to_picture(slide, 6, hp_fig_path)
    
            mr_fig = mret_hmap(b)
            mr_fig_path = save_path + '/mr_fig.png'
            mr_fig.savefig(mr_fig_path, bbox_inches = 'tight')
            chart_to_picture(slide, 15, mr_fig_path)
        
        else:
            pass
        
        
      #----------------------------------------------------------------------------
      # bottom part of a slide
        
        if ticker in Equity + Real_Estate + Multi_Asset + New:
           # 1. Trading Volume   
            text_to_text_placeholder(slide, 7, '월간 평균 거래량(천 주)')
            v = make_tradibility(soup1)
            vp_fig = v
            vp_fig_path = save_path + '/vp_fig.png'
            vp_fig.savefig(vp_fig_path, bbox_inches = 'tight')
            chart_to_picture(slide, 8, vp_fig_path)        


        # 3. Holding Analysis   
            text_to_text_placeholder(slide, 11, '섹터 분포(%)')
            hap = holding_analysis(soup2)
            hap_fig = hap
            hap_fig_path = save_path + '/hap_fig.png'
            hap_fig.savefig(hap_fig_path, dpi = 'figure', bbox_inches = 'tight')
            chart_to_picture(slide, 12, hap_fig_path)   


        elif ticker in Equity_Global + Commodity + Currency:
           # 1. Trading Volume   
            text_to_text_placeholder(slide, 7, '월간 평균 거래량(천 주)')
            v = make_tradibility(soup1)
            vp_fig = v
            vp_fig_path = save_path + '/vp_fig.png'
            vp_fig.savefig(vp_fig_path, bbox_inches = 'tight')
            chart_to_picture(slide, 8, vp_fig_path)        
       
        # 2. Tracking Error   
            text_to_text_placeholder(slide, 11, '추적 오차(%)')
            t = make_tracking_error(soup1)
            tp_fig = t
            tp_fig_path = save_path + '/tp_fig.png'
            tp_fig.savefig(tp_fig_path, dpi = 'figure', bbox_inches = 'tight')
            chart_to_picture(slide, 12, tp_fig_path)   
        
        elif ticker in Fixed_Income:
           # 1. Trading Volume   
            text_to_text_placeholder(slide, 7, '월간 평균 거래량(천 주)')
            v = make_tradibility(soup1)
            vp_fig = v
            vp_fig_path = save_path + '/vp_fig.png'
            vp_fig.savefig(vp_fig_path, bbox_inches = 'tight')
            chart_to_picture(slide, 8, vp_fig_path)        


        # 3. Holding Analysis   
            text_to_text_placeholder(slide, 11, '섹터 분포(%)')
            bhap = bond_holding_analysis(soup2)
            bhap_fig = bhap
            bhap_fig_path = save_path + '/bhap_fig.png'
            bhap_fig.savefig(bhap_fig_path, dpi = 'figure', bbox_inches = 'tight')
            chart_to_picture(slide, 12, bhap_fig_path)   

        else:
            pass
        
        
        # Classification
        # Equity
        if ticker in Equity_KR:
            cl = '주식 > 한국'
        elif ticker in Equity_KR_Sec:
            cl = '주식 > 한국 섹터'
        elif ticker in Equity_KR_SM:
            cl = '주식 > 한국 스마트베타'
        elif ticker in Equity_KR_TM:
            cl = '주식 > 한국 테마형'
        elif ticker in Equity_Glb:
            cl = '주식 > 글로벌 주식'
        elif ticker in Equity_Glb_TM:
            cl = '주식 > 글로벌 테마형'
        elif ticker in Equity_US:
            cl = '주식 > 미국'
        elif ticker in Equity_US_TM:
            cl = '주식 > 미국 테마형'
        elif ticker in Equity_US_Sec:
            cl = '주식 > 미국 섹터'
        elif ticker in Equity_EU:
            cl = '주식 > 유럽'
        elif ticker in Equity_JP:
            cl = '주식 > 일본'
        elif ticker in Equity_CH:
            cl = '주식 > 중국'
        elif ticker in Equity_CH_TM:
            cl = '주식 > 중국 테마형'            
  
            
            
        #
        # Fixed Income
        elif ticker in FI_Govt:
            cl = '채권 > 국채'
        elif ticker in FI_US:
            cl = '채권 > 미국'
        elif ticker in FI_AGG:
            cl = '채권 > 종합'
        elif ticker in FI_Credit:
            cl = '채권 > 크레딧'
        elif ticker in FI_Others:
            cl = '채권 > 기타'
        #
        # Commodity
        elif ticker in Comm_Energy:
            cl = '원자재 > 에너지'
        elif ticker in Comm_Precious:
            cl = '원자재 > 귀금속'
        elif ticker in Comm_Agri:
            cl = '원자재 > 농산물'
        #
        # Real Estate
        elif ticker in RE_AGG:
            cl = '리츠 > 종합'
        #
        # Currency
        elif ticker in Currency:
            cl = '환율'
        #
        # Multi Asset
        elif ticker in Multi_Asset:
            cl = '멀티에셋/자산배분'
        
        #
        # Uncategorized    
        elif ticker in New:
            cl = '신규편입'    

    
        text_to_text_placeholder(slide, 16, cl + ' > ' + make_index(soup2))        
        
        prs.save(save_path + '/etf_kd_ver1.pptx')
        print("%s: complete" %ticker)
        plt.clf()
    except:
        print('%s: Fail' %ticker)

# %% Classification

Equity_KR = ['A100910', 'A292190', 'A229720', 'A156080', 'A229200', 'A226490', 'A069500', 'A278530']
Equity_KR_Sec = ['A139260', 'A326240', 'A117700', 'A266390', 'A139270', 'A102960', 'A244580', 'A091160', 'A098560', 'A140700', 'A227550',\
                 'A227560', 'A157490', 'A139250', 'A228800', 'A140710', 'A091170', 'A307510', 'A091180', 'A139230',\
                 'A102970', 'A117680', 'A143860', 'A139280', 'A315270']
Equity_KR_SM = ['A272230', 'A223190', 'A333960', 'A252730', 'A279540', 'A333980', 'A161510', 'A315960']
Equity_KR_TM = ['A305720', 'A322400', 'A138530', 'A292150', 'A315930', 'A326230', 'A314700', 'A228810', 'A266360',\
                'A261070', 'A091230', 'A102780', 'A325010', 'A150460', 'A138540', 'A228790', 'A290130', 'A278420', 'A285690',\
                'A376410', 'A377990', 'A368680', 'A367770', 'A364980', 'A364990', 'A364970', 'A365000', 'A385720', 'A400970']
                
                
Equity_Glb = ['A251350', 'A195980']
Equity_Glb_TM = ['A275980', 'A276990', 'A298770', 'A269420', 'A276000', 'A248270', 'A387270',\
                 'A394660', 'A354350', 'A391590', 'A391600']
Equity_US = ['A245340', 'A280930', 'A133690', 'A143850']
Equity_US_TM = ['A314250', 'A309230', 'A276970', 'A287180', 'A381170', 'A381180', 'A200030']
Equity_US_Sec = ['A280320', 'A203780']
Equity_EU = ['A195930', 'A245350']
Equity_JP = ['A241180']
Equity_CH = ['A174360', 'A192090', 'A169950', 'A245360', 'A310080', 'A256750'] 
Equity_CH_TM = ['A371470', 'A371460', 'A371150']
    
Equity = Equity_KR + Equity_KR_Sec + Equity_KR_SM + Equity_KR_TM

Equity_Global = Equity_Glb + Equity_Glb_TM + Equity_US +\
                Equity_US_TM + Equity_US_Sec + Equity_EU + Equity_JP + Equity_CH + Equity_CH_TM

FI_Govt = ['A272910', 'A148070', 'A272560', 'A114260']
FI_US = ['A305080', 'A304660', 'A329750']
FI_AGG = ['A153130', 'A214980', 'A273130']
FI_Credit = ['A136340', 'A332610', 'A332620', 'A239660', 'A182490']
FI_Others = ['A336160', 'A273140', 'A196230']

Fixed_Income = FI_Govt + FI_US + FI_AGG + FI_Credit + FI_Others

Comm_Energy = ['A261220', 'A400570']
Comm_Precious = ['A138910', 'A132030', 'A144600', 'A334690']
Comm_Agri = ['A137610']
Commodity = Comm_Energy + Comm_Precious + Comm_Agri

RE_AGG = ['A329200', 'A352560', 'A316300']
Real_Estate = RE_AGG

Multi_Asset = ['A183710', 'A253290', 'A237440', 'A183700',  'A237370', 'A251600',\
               'A284430', 'A321410', 'A241390']

Currency = ['A261240']


univ = Equity + Equity_Global + Fixed_Income + Commodity + Real_Estate + Multi_Asset + Currency




# %% Main

for i in Equity:
    ticker = i
    make_slide(ticker)
    
for i in Equity_Global:
    ticker = i
    make_slide(ticker)    

for i in Fixed_Income:
    ticker = i
    make_slide(ticker)

for i in Commodity:
    ticker = i
    make_slide(ticker)
    
for i in Real_Estate:
    ticker = i
    make_slide(ticker)

for i in Multi_Asset:
    ticker = i
    make_slide(ticker)
    
for i in Currency:
    ticker = i
    make_slide(ticker)  

for i in New:
    ticker = i
    make_slide(ticker)
