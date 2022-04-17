from pptx import Presentation


    
def chart_to_picture(slide, placeholder_number, fig_path):
    picture_placeholder = slide.shapes[placeholder_number]
    picture_placeholder.insert_picture(fig_path)

def chart_to_picture_style(slide, placeholder_number, fig_path):
    slide.placeholders[placeholder_number].insert_picture(fig_path)


save_path = r'C:/ETF'

prs = Presentation(save_path + '/etf_gd_ver2.pptx')
for idx, slide_layout in enumerate(prs.slide_layouts):
    slide = prs.slides.add_slide(slide_layout)
    for shape in slide.placeholders:
        print(idx, shape.placeholder_format.idx, shape.name, shape.top, shape.left)

prs = Presentation(save_path + '/etf_gd_ver2.pptx')    # open a pptx file
layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(layout)

pf_dict = {'SPY Portfolio Data ': {'Weighted Average Market Cap': '$593.42B ', 'Price / Earnings Ratio': 
'21.73 ', 'Price / Book Ratio': '4.32 ', 'Distribution Yield': '1.29% ', 'Next Ex-Dividend Date': '06/17/22 ', 'Number of Holdings': '501 '}}

pf_values = pf_dict.values()
print(type(pf_values))
print(pf_values)
for k,v in pf_values.__iter__:
    print(k, v)